from __future__ import annotations

from pathlib import Path
import sys
from pptx import Presentation


def slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        if hasattr(sh, "text") and sh.text:
            parts.append(sh.text)
    return "\n".join(parts)


def titleish(text: str) -> str:
    for line in text.splitlines():
        line = line.strip()
        if line:
            return line[:80]
    return ""


def main() -> int:
    pptx_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("report/auto_deck.pptx")
    if not pptx_path.exists():
        print(f"ERROR: missing deck: {pptx_path}")
        return 2

    p = Presentation(str(pptx_path))
    missing = []
    for i, slide in enumerate(p.slides, 1):
        txt = slide_text(slide)
        if "How to interpret" not in txt:
            missing.append((i, titleish(txt)))

    print(f"slides={len(p.slides)}")
    print(f"slides_missing_interpret={len(missing)}")
    if missing:
        print("missing (index, title-ish):")
        for idx, t in missing[:50]:
            print(f"  {idx:02d}: {t}")
        return 2

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
