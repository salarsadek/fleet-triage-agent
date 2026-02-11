"""
Professional auto-deck generator (PowerPoint + PDF).

Adds:
- "How to interpret" panel on core slides
- "Key takeaways" slide per major section (Data quality / EDA+Stats / Model performance / Fleet triage)
- Keeps appendix (extra figures/tables)

Usage:
  python -m src.reporting.make_auto_deck --config config.yaml --out report/auto_deck.pptx
"""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional, Tuple

import pandas as pd
import yaml

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

# Optional for perfect aspect-ratio fitting. Fallback if not installed.
try:
    from PIL import Image  # type: ignore
except Exception:  # pragma: no cover
    Image = None  # type: ignore


# -----------------------------
# Visual style
# -----------------------------
SLIDE_W = Inches(13.333)   # 16:9 widescreen
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.65)
MARGIN_R = Inches(0.65)

HEADER_H = Inches(0.72)
FOOTER_H = Inches(0.25)

BODY_TOP = HEADER_H + Inches(0.40)
BODY_BOTTOM = SLIDE_H - FOOTER_H - Inches(0.12)
BODY_H = BODY_BOTTOM - BODY_TOP

GAP = Inches(0.35)
PAD = Inches(0.18)

NOTE_H = Inches(1.55)  # interpretation panel height
MEDIA_H = BODY_H - NOTE_H - Inches(0.18)

ACCENT = RGBColor(0x1F, 0x4E, 0x79)      # deep blue
ACCENT_2 = RGBColor(0x12, 0x2B, 0x45)    # darker blue
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)

CARD_BG = RGBColor(0xF7, 0xF8, 0xFA)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


@dataclass
class Paths:
    root: Path
    tables: Path
    figures: Path
    reports: Path


def repo_root() -> Path:
    return Path(__file__).resolve().parents[2]


def load_yaml(path: Path) -> dict:
    # robust against BOM in config.yaml
    with open(path, "r", encoding="utf-8-sig") as f:
        cfg = yaml.safe_load(f)
    if not isinstance(cfg, dict):
        raise ValueError("Config did not parse as a dictionary.")
    return cfg


def ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


def safe_read_csv(path: Path) -> Optional[pd.DataFrame]:
    if not path.exists():
        return None
    try:
        return pd.read_csv(path)
    except Exception:
        return None


def safe_read_json(path: Path) -> Optional[dict]:
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None


def fmt_val(x) -> str:
    try:
        if pd.isna(x):
            return ""
    except Exception:
        pass

    if isinstance(x, float):
        if abs(x) >= 1000:
            return f"{x:,.0f}"
        if abs(x) >= 10:
            return f"{x:.2f}"
        return f"{x:.3f}"
    return str(x)


def infer_stamp_model_horizon(paths: Paths) -> Tuple[str, str, str]:
    stamp_path = paths.reports / "triage_stamp_latest.txt"
    stamp = stamp_path.read_text(encoding="utf-8").strip() if stamp_path.exists() else "latest"

    parts = stamp.split("_")
    horizon = parts[1] if len(parts) >= 2 else "30d"
    model = parts[2] if len(parts) >= 3 else "hgb"
    return stamp, model, horizon


# -----------------------------
# PPT helpers
# -----------------------------
def set_widescreen(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_header(slide, title: str, subtitle: Optional[str] = None) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    tx = slide.shapes.add_textbox(MARGIN_L, Inches(0.12), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.40))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        tx2 = slide.shapes.add_textbox(MARGIN_L, Inches(0.48), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.28))
        tf2 = tx2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = FONT_BODY
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF5)


def add_footer(slide, left_text: str, right_text: str) -> None:
    y = SLIDE_H - FOOTER_H

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y - Inches(0.02), SLIDE_W, Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = CARD_BORDER
    ln.line.fill.background()

    lx = slide.shapes.add_textbox(MARGIN_L, y, SLIDE_W / 2, FOOTER_H)
    ltf = lx.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lp.text = left_text
    lp.font.name = FONT_BODY
    lp.font.size = Pt(10)
    lp.font.color.rgb = TEXT_MUTED

    rx = slide.shapes.add_textbox(SLIDE_W / 2, y, SLIDE_W / 2 - MARGIN_R, FOOTER_H)
    rtf = rx.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.text = right_text
    rp.alignment = PP_ALIGN.RIGHT
    rp.font.name = FONT_BODY
    rp.font.size = Pt(10)
    rp.font.color.rgb = TEXT_MUTED


def slide_blank(prs: Presentation, title: str, subtitle: Optional[str], footer_left: str, footer_right: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)
    add_footer(slide, footer_left, footer_right)
    return slide


def add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    return card


def _fit_image_size(img_path: Path, box_w, box_h) -> Tuple[int, int]:
    if Image is None:
        return int(box_w), 0  # width-only fallback

    try:
        with Image.open(img_path) as im:
            w_px, h_px = im.size
        if w_px <= 0 or h_px <= 0:
            return int(box_w), 0

        img_aspect = w_px / h_px
        box_aspect = float(box_w) / float(box_h)

        if img_aspect >= box_aspect:
            w = float(box_w)
            h = w / img_aspect
        else:
            h = float(box_h)
            w = h * img_aspect

        return int(w), int(h)
    except Exception:
        return int(box_w), 0


def add_picture_fit(slide, img: Path, left, top, box_w, box_h):
    w, h = _fit_image_size(img, box_w, box_h)
    if h == 0:
        slide.shapes.add_picture(str(img), left, top, width=box_w)
        return
    x = left + int((box_w - w) / 2)
    y = top + int((box_h - h) / 2)
    slide.shapes.add_picture(str(img), x, y, width=w, height=h)


def add_media_card(slide, img: Path, left, top, width, height, caption: Optional[str] = None):
    add_card(slide, left, top, width, height)

    cap_h = Inches(0.34) if caption else Inches(0.0)
    inner_left = left + PAD
    inner_top = top + PAD
    inner_w = width - 2 * PAD
    inner_h = height - 2 * PAD - cap_h

    add_picture_fit(slide, img, inner_left, inner_top, inner_w, inner_h)

    if caption:
        tx = slide.shapes.add_textbox(inner_left, top + height - cap_h - PAD / 2, inner_w, cap_h)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = caption
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED


def add_bullets(slide, bullets: Iterable[str], left, top, width, height, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    bullets = list(bullets) or ["(no content)"]
    p0 = tf.paragraphs[0]
    p0.text = f"• {bullets[0]}"
    p0.font.name = FONT_BODY
    p0.font.size = Pt(font_size)
    p0.font.color.rgb = TEXT_DARK

    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK


def add_interpretation_panel(slide, bullets: list[str], left, top, width, height):
    # Bottom panel: "How to interpret"
    add_card(slide, left, top, width, height)

    # Title
    tx = slide.shapes.add_textbox(left + PAD, top + PAD / 2, width - 2 * PAD, Inches(0.30))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "How to interpret"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_2

    # Bullets
    add_bullets(
        slide,
        bullets,
        left=left + PAD,
        top=top + Inches(0.35),
        width=width - 2 * PAD,
        height=height - Inches(0.40),
        font_size=13,
    )


def add_section_divider(prs: Presentation, footer_left: str, footer_right: str, title: str, subtitle: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_2
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(MARGIN_L, Inches(2.6), SLIDE_W - MARGIN_L - MARGIN_R, Inches(1.0))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        tx2 = slide.shapes.add_textbox(MARGIN_L, Inches(3.5), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.6))
        tf2 = tx2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = FONT_BODY
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xDD, 0xE7, 0xF3)

    add_footer(slide, footer_left, footer_right)
    return slide


def add_kpi_tiles(slide, items: list[Tuple[str, str]], left, top, width, height):
    n = min(len(items), 4)
    if n == 0:
        return
    gap = Inches(0.25)
    tile_w = (width - gap * (n - 1)) / n
    tile_h = height

    for i in range(n):
        label, value = items[i]
        x = left + i * (tile_w + gap)
        add_card(slide, x, top, tile_w, tile_h)

        txv = slide.shapes.add_textbox(x + PAD, top + Inches(0.22), tile_w - 2 * PAD, Inches(0.60))
        tfv = txv.text_frame
        tfv.clear()
        pv = tfv.paragraphs[0]
        pv.text = value
        pv.font.name = FONT_TITLE
        pv.font.size = Pt(28)
        pv.font.bold = True
        pv.font.color.rgb = TEXT_DARK

        txl = slide.shapes.add_textbox(x + PAD, top + Inches(0.90), tile_w - 2 * PAD, Inches(0.35))
        tfl = txl.text_frame
        tfl.clear()
        pl = tfl.paragraphs[0]
        pl.text = label
        pl.font.name = FONT_BODY
        pl.font.size = Pt(12)
        pl.font.color.rgb = TEXT_MUTED


# -----------------------------
# Table rendering (CSV -> PNG)
# -----------------------------
def render_table_png(
    df: pd.DataFrame,
    out_path: Path,
    title: Optional[str] = None,
    max_rows: int = 18,
    max_cols: int = 12,
) -> Tuple[Path, bool]:
    df2 = df.copy()

    if df2.shape[1] > max_cols:
        df2 = df2.iloc[:, :max_cols]

    truncated = False
    if df2.shape[0] > max_rows:
        df2 = df2.iloc[:max_rows, :]
        truncated = True

    # elementwise formatting without FutureWarning
    if hasattr(df2, "map"):
        df2 = df2.map(fmt_val)  # pandas >= 2.1+
    else:  # pragma: no cover
        df2 = df2.applymap(fmt_val)

    rows, _cols = df2.shape
    fig_h = min(0.36 * (rows + 1) + (0.6 if title else 0.2), 6.0)
    fig_w = 12.0

    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=240)
    ax.axis("off")

    if title:
        ax.set_title(title, fontsize=14, fontweight="bold", loc="left", pad=6)

    tbl = ax.table(
        cellText=df2.values,
        colLabels=[str(c) for c in df2.columns],
        cellLoc="left",
        colLoc="left",
        loc="upper left",
    )

    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9.5)
    tbl.scale(1.0, 1.18)

    for (r, _c), cell in tbl.get_celld().items():
        cell.set_linewidth(0.3)
        cell.set_edgecolor("#DDDDDD")
        if r == 0:
            cell.set_facecolor("#F2F4F7")
            cell.set_text_props(weight="bold", color="#111111")
        else:
            cell.set_facecolor("#FFFFFF")

    if truncated:
        ax.text(0.0, -0.05, f"Showing first {max_rows} rows (truncated).", fontsize=9, color="#666666", transform=ax.transAxes)

    ensure_dir(out_path.parent)
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)

    return out_path, truncated


# -----------------------------
# Small extractors for takeaways
# -----------------------------
def _pick_col(df: pd.DataFrame, patterns: list[str]) -> Optional[str]:
    cols = [c for c in df.columns]
    low = {c: str(c).lower() for c in cols}
    for pat in patterns:
        rgx = re.compile(pat)
        for c in cols:
            if rgx.search(low[c]):
                return c
    return None


def summarize_work_orders(df: pd.DataFrame, topn: int = 2) -> list[str]:
    if df is None or df.empty:
        return []
    # guess columns
    c_sub = _pick_col(df, [r"subsystem", r"system", r"category", r"type"])
    c_cnt = _pick_col(df, [r"count", r"n$", r"num", r"total"])
    if c_sub is None:
        c_sub = df.columns[0]
    if c_cnt is None and len(df.columns) >= 2:
        c_cnt = df.columns[1]
    if c_cnt is None:
        return []
    d2 = df[[c_sub, c_cnt]].copy()
    try:
        d2[c_cnt] = pd.to_numeric(d2[c_cnt], errors="coerce")
        d2 = d2.dropna(subset=[c_cnt]).sort_values(c_cnt, ascending=False)
    except Exception:
        return []
    items = []
    for _, r in d2.head(topn).iterrows():
        items.append(f"Highest maintenance demand: {r[c_sub]} ({int(r[c_cnt])} work orders).")
    return items


def summarize_cox(df: pd.DataFrame, topn: int = 2) -> list[str]:
    if df is None or df.empty:
        return []
    c_feat = _pick_col(df, [r"feature", r"covariate", r"variable", r"name"])
    c_hr = _pick_col(df, [r"hazard.*ratio", r"\bhr\b"])
    if c_feat is None:
        c_feat = df.columns[0]
    if c_hr is None:
        # try second column
        if len(df.columns) >= 2:
            c_hr = df.columns[1]
        else:
            return []
    d2 = df[[c_feat, c_hr]].copy()
    try:
        d2[c_hr] = pd.to_numeric(d2[c_hr], errors="coerce")
        d2 = d2.dropna(subset=[c_hr]).sort_values(c_hr, ascending=False)
    except Exception:
        return []
    items = []
    for _, r in d2.head(topn).iterrows():
        hr = float(r[c_hr])
        if hr >= 1.0:
            items.append(f"Top risk-increasing factor: {r[c_feat]} (HR≈{hr:.2f}; >1 increases hazard).")
    return items


def summarize_risk_metrics(df: pd.DataFrame) -> list[str]:
    if df is None or df.empty:
        return []
    c_auprc = _pick_col(df, [r"auprc"])
    if c_auprc is None:
        return []
    d2 = df.copy()
    try:
        d2[c_auprc] = pd.to_numeric(d2[c_auprc], errors="coerce")
        d2 = d2.dropna(subset=[c_auprc]).sort_values(c_auprc, ascending=False)
    except Exception:
        return []
    best = d2.iloc[0]
    model = best.get("model", "?")
    task = best.get("task", "?")
    auprc = float(best[c_auprc])
    out = [f"Best ranking quality (AUPRC): {auprc:.3f} for {task} / {model}."]
    c_brier = _pick_col(d2, [r"brier"])
    if c_brier is not None and pd.notna(best.get(c_brier)):
        try:
            out.append(f"Lower Brier is better (probabilities). Best Brier: {float(best[c_brier]):.3f}.")
        except Exception:
            pass
    return out


def find_first(paths: list[Path]) -> Optional[Path]:
    for p in paths:
        if p.exists():
            return p
    return None


def is_table_candidate(p: Path) -> bool:
    n = p.name.lower()
    if not n.endswith(".csv"):
        return False
    if n.startswith("test_predictions_"):
        return False
    return True


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser()
    p.add_argument("--config", type=str, required=True)
    p.add_argument("--out", type=str, default="report/auto_deck.pptx")
    p.add_argument("--max_appendix_figures", type=int, default=40)
    p.add_argument("--max_appendix_tables", type=int, default=20)
    return p.parse_args()


def main() -> None:
    args = parse_args()
    root = repo_root()

    cfg_path = Path(args.config)
    if not cfg_path.is_absolute():
        cfg_path = root / cfg_path
    cfg = load_yaml(cfg_path)

    paths = Paths(
        root=root,
        tables=root / cfg["paths"]["tables_dir"],
        figures=root / cfg["paths"]["figures_dir"],
        reports=root / cfg["paths"]["reports_dir"],
    )

    stamp, model, horizon = infer_stamp_model_horizon(paths)

    footer_left = "fleet-triage-agent • auto-generated deck"
    footer_right = f"snapshot: {stamp}"

    prs = Presentation()
    set_widescreen(prs)

    deck_assets = root / "report" / "_deck_assets"
    ensure_dir(deck_assets)

    used_figs: set[str] = set()
    used_tables: set[str] = set()
    missing: list[str] = []

    full_w = SLIDE_W - MARGIN_L - MARGIN_R
    col_w = (full_w - GAP) / 2
    note_y = BODY_TOP + MEDIA_H + Inches(0.18)

    # -----------------------------
    # Title
    # -----------------------------
    slide = slide_blank(prs, "Fleet Maintenance Triage Agent", "Results deck", footer_left, footer_right)
    add_bullets(
        slide,
        [
            "Objective: prioritize vehicles for service using predicted near-term breakdown risk and operational cost.",
            "Reproducible pipeline: validate → EDA/stats → train → triage snapshot → latest aliases → deck/PDF.",
            "Slides are built from exported artifacts in outputs/ (figures + tables).",
        ],
        left=MARGIN_L, top=BODY_TOP, width=full_w, height=Inches(2.2), font_size=18
    )
    add_interpretation_panel(
        slide,
        [
            "This deck is fully reproducible: rerun the pipeline to refresh numbers and visuals.",
            "The ‘snapshot’ in the footer tells you which triage run the deck is based on.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Agenda
    slide = slide_blank(prs, "Agenda", None, footer_left, footer_right)
    add_bullets(
        slide,
        [
            "Data quality & dataset overview",
            "EDA + reliability/statistics (KM + Cox)",
            "Model performance (PR + calibration) + metrics",
            "Fleet triage view + recommended queue",
            "Guardrails + similar-case evidence",
            "Appendix: additional artifacts",
        ],
        left=MARGIN_L, top=BODY_TOP, width=full_w, height=MEDIA_H, font_size=18
    )
    add_interpretation_panel(
        slide,
        [
            "Main slides tell the story end-to-end.",
            "Appendix contains extra plots/tables auto-included from outputs/.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Pipeline
    slide = slide_blank(prs, "Pipeline", "One-command reproducible workflow", footer_left, footer_right)
    add_bullets(
        slide,
        [
            r"Run full pipeline:  .\scripts\run.ps1 report",
            r"Or rebuild deck only: .\scripts\run.ps1 deck",
        ],
        left=MARGIN_L, top=BODY_TOP, width=full_w, height=Inches(1.2), font_size=16
    )
    add_interpretation_panel(
        slide,
        [
            "Validate catches data issues early (schema, ranges, missingness).",
            "EDA/Stats gives context and risk drivers (KM + Cox).",
            "Train exports PR/calibration + metrics; Triage exports queue + guardrails + evidence.",
        ],
        left=MARGIN_L, top=BODY_TOP + Inches(1.45), width=full_w, height=BODY_H - Inches(1.45)
    )

    # -----------------------------
    # Executive summary
    # -----------------------------
    slide = slide_blank(prs, "Executive summary", None, footer_left, footer_right)

    rm = safe_read_csv(paths.tables / "risk_metrics.csv")
    best_auprc = None
    best_brier = None
    if rm is not None and not rm.empty and "auprc" in rm.columns:
        try:
            rm2 = rm.copy()
            rm2["auprc"] = pd.to_numeric(rm2["auprc"], errors="coerce")
            rm2 = rm2.dropna(subset=["auprc"]).sort_values("auprc", ascending=False)
            best = rm2.iloc[0]
            best_auprc = float(best.get("auprc", float("nan")))
            if "brier" in rm2.columns:
                best_brier = float(best.get("brier", float("nan")))
            used_tables.add("risk_metrics.csv")
        except Exception:
            pass

    triage_sum = safe_read_json(paths.tables / "triage_summary_latest.json")
    actionable = triage_sum.get("n_actionable_ok") if isinstance(triage_sum, dict) else None
    abstain = triage_sum.get("n_abstain") if isinstance(triage_sum, dict) else None

    tiles = []
    if best_auprc is not None and best_auprc == best_auprc:
        tiles.append(("Best AUPRC", f"{best_auprc:.3f}"))
    if best_brier is not None and best_brier == best_brier:
        tiles.append(("Best Brier", f"{best_brier:.3f}"))
    tiles.append(("Model / Horizon", f"{model} / {horizon}"))
    if actionable is not None and abstain is not None:
        tiles.append(("Actionable / Abstain", f"{actionable} / {abstain}"))

    add_kpi_tiles(slide, tiles, MARGIN_L, BODY_TOP, full_w, Inches(1.45))
    add_interpretation_panel(
        slide,
        [
            "AUPRC is best for rare events (higher is better).",
            "Brier summarizes probability quality (lower is better).",
            "Actionable vs Abstain shows how often guardrails trigger manual review.",
        ],
        left=MARGIN_L, top=BODY_TOP + Inches(1.65), width=full_w, height=BODY_H - Inches(1.65)
    )

    # -----------------------------
    # Data quality
    # -----------------------------
    add_section_divider(prs, footer_left, footer_right, "Data quality", "Validation + dataset overview")

    slide = slide_blank(prs, "Data quality", "Validation summary", footer_left, footer_right)
    vdf = safe_read_csv(paths.tables / "validation_summary_table.csv")
    if vdf is not None:
        img_path, _ = render_table_png(vdf, deck_assets / "validation_summary.png", title="Validation summary", max_rows=22, max_cols=6)
        add_media_card(slide, img_path, MARGIN_L, BODY_TOP, full_w, MEDIA_H, caption="Great Expectations checks")
        used_tables.add("validation_summary_table.csv")
    else:
        missing.append("outputs/tables/validation_summary_table.csv")
        add_bullets(slide, ["Missing validation summary. Run: .\\scripts\\run.ps1 validate"], MARGIN_L, BODY_TOP, full_w, MEDIA_H)

    add_interpretation_panel(
        slide,
        [
            "Each row is a data quality rule (schema/range/missingness).",
            "All checks should pass in normal operation; failures mean unreliable downstream results.",
            "Use this slide to justify that the pipeline is safe to run automatically.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Key takeaways (Data quality)
    slide = slide_blank(prs, "Key takeaways", "Data quality", footer_left, footer_right)
    vrep = safe_read_json(paths.tables / "validation_report.json")
    bullets = [
        "Goal: ensure inputs are consistent (types, ranges, missingness) before modeling.",
    ]
    if isinstance(vrep, dict) and vrep.get("checks_passed") is not None and vrep.get("checks_total") is not None:
        bullets.append(f"Checks passed: {vrep['checks_passed']}/{vrep['checks_total']} (pipeline gate).")
        used_tables.add("validation_report.json")
    bullets.extend(
        [
            "If checks fail, rerun data generation / fix upstream ingestion before trusting any metrics.",
            "This is the foundation for a production-quality ML/agent workflow.",
        ]
    )
    add_bullets(slide, bullets, MARGIN_L, BODY_TOP, full_w, BODY_H, font_size=18)

    # -----------------------------
    # EDA + stats
    # -----------------------------
    add_section_divider(prs, footer_left, footer_right, "EDA & statistics", "Operational patterns and risk factors")

    slide = slide_blank(prs, "EDA overview", "Dataset + operational signals", footer_left, footer_right)
    eda_df = safe_read_csv(paths.tables / "eda_overview_table.csv")
    wo_df = safe_read_csv(paths.tables / "work_orders_by_subsystem.csv")

    if eda_df is not None:
        img1, _ = render_table_png(eda_df, deck_assets / "eda_overview.png", title="EDA overview", max_rows=16, max_cols=10)
        add_media_card(slide, img1, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Dataset summary")
        used_tables.add("eda_overview_table.csv")
    else:
        missing.append("outputs/tables/eda_overview_table.csv")

    if wo_df is not None:
        img2, _ = render_table_png(wo_df, deck_assets / "work_orders.png", title="Work orders by subsystem", max_rows=16, max_cols=8)
        add_media_card(slide, img2, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Maintenance demand")
        used_tables.add("work_orders_by_subsystem.csv")
    else:
        missing.append("outputs/tables/work_orders_by_subsystem.csv")

    add_interpretation_panel(
        slide,
        [
            "EDA overview: sanity-check volume, coverage, and key variables (are values plausible?).",
            "Work orders: higher counts indicate where maintenance load is concentrated.",
            "Use this to connect the model to real operations (what breaks, where, and how often).",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    slide = slide_blank(prs, "Reliability & risk factors", "Survival + hazard ratios", footer_left, footer_right)
    km = paths.figures / "km_duty_cycle.png"
    cox = safe_read_csv(paths.tables / "cox_hazard_ratios.csv")

    if km.exists():
        add_media_card(slide, km, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Kaplan–Meier by duty-cycle group")
        used_figs.add(km.name)
    else:
        missing.append("outputs/figures/km_duty_cycle.png")

    if cox is not None:
        img, _ = render_table_png(cox, deck_assets / "cox_hazards.png", title="Cox hazard ratios", max_rows=16, max_cols=8)
        add_media_card(slide, img, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Risk multipliers (HR)")
        used_tables.add("cox_hazard_ratios.csv")
    else:
        missing.append("outputs/tables/cox_hazard_ratios.csv")

    add_interpretation_panel(
        slide,
        [
            "Kaplan–Meier: lower curve means faster failures; separation suggests duty cycle impacts reliability.",
            "Cox hazard ratios: HR>1 increases failure hazard, HR<1 decreases; bigger deviation from 1 = stronger effect.",
            "These results help explain *why* certain vehicles get prioritized (not only that they do).",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Key takeaways (EDA & stats)
    slide = slide_blank(prs, "Key takeaways", "EDA & statistics", footer_left, footer_right)
    bullets = [
        "EDA connects model outputs to fleet reality: coverage, usage patterns, and failure/service signals.",
    ]
    if wo_df is not None:
        bullets.extend(summarize_work_orders(wo_df, topn=2))
    else:
        bullets.append("Work orders: missing table → rerun report to regenerate.")
    if cox is not None:
        bullets.extend(summarize_cox(cox, topn=2))
    else:
        bullets.append("Cox model: missing hazard ratios → rerun report to regenerate.")
    bullets.append("Use these insights to propose targeted interventions (routes, duty cycles, subsystems).")
    add_bullets(slide, bullets, MARGIN_L, BODY_TOP, full_w, BODY_H, font_size=18)

    # -----------------------------
    # Model performance
    # -----------------------------
    add_section_divider(prs, footer_left, footer_right, "Model performance", "Calibration + ranking quality")

    slide = slide_blank(prs, "Model performance", f"{model} • {horizon}", footer_left, footer_right)
    pr = paths.figures / f"pr_{model}_{horizon}.png"
    cal = paths.figures / f"calibration_{model}_{horizon}.png"

    if pr.exists():
        add_media_card(slide, pr, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Precision–Recall")
        used_figs.add(pr.name)
    else:
        missing.append(f"outputs/figures/pr_{model}_{horizon}.png")

    if cal.exists():
        add_media_card(slide, cal, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Calibration")
        used_figs.add(cal.name)
    else:
        missing.append(f"outputs/figures/calibration_{model}_{horizon}.png")

    add_interpretation_panel(
        slide,
        [
            "Precision–Recall: higher is better; focus on precision at the recall region you need operationally.",
            "Calibration: points close to the diagonal mean predicted probabilities match real-world frequencies.",
            "Good ranking (PR) helps choose *which* vehicles; good calibration helps decide *how urgent* they are.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    slide = slide_blank(prs, "Metrics summary", "Model ranking + operational precision", footer_left, footer_right)
    rm_df = safe_read_csv(paths.tables / "risk_metrics.csv")
    pk_df = safe_read_csv(paths.tables / "precision_at_k_mondays.csv")

    if rm_df is not None:
        if "auprc" in rm_df.columns:
            try:
                rm_df = rm_df.copy()
                rm_df["auprc"] = pd.to_numeric(rm_df["auprc"], errors="coerce")
                rm_df = rm_df.sort_values("auprc", ascending=False)
            except Exception:
                pass
        img, _ = render_table_png(rm_df, deck_assets / "risk_metrics.png", title="Risk metrics (sorted)", max_rows=12, max_cols=12)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Overall ranking")
        used_tables.add("risk_metrics.csv")
    else:
        missing.append("outputs/tables/risk_metrics.csv")

    if pk_df is not None:
        img, _ = render_table_png(pk_df, deck_assets / "precision_at_k.png", title="Precision@K (Mondays)", max_rows=14, max_cols=10)
        add_media_card(slide, img, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Top-K quality")
        used_tables.add("precision_at_k_mondays.csv")
    else:
        missing.append("outputs/tables/precision_at_k_mondays.csv")

    add_interpretation_panel(
        slide,
        [
            "Risk metrics: compare models; higher AUPRC is better for rare events.",
            "Precision@K: if K is your service capacity, this approximates how many queued vehicles are truly high-risk.",
            "Use these tables to justify the selected model and why it’s operationally useful.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Key takeaways (Model performance)
    slide = slide_blank(prs, "Key takeaways", "Model performance", footer_left, footer_right)
    bullets = [
        "We care about both: (1) ranking the right vehicles, and (2) reliable probabilities for urgency.",
    ]
    if rm_df is not None:
        bullets.extend(summarize_risk_metrics(rm_df))
    else:
        bullets.append("Risk metrics missing → rerun report to regenerate.")
    bullets.extend(
        [
            "PR supports ‘who to service first’; calibration supports ‘how risky is this vehicle’.",
            "Precision@K maps model performance to a real service queue size.",
        ]
    )
    add_bullets(slide, bullets, MARGIN_L, BODY_TOP, full_w, BODY_H, font_size=18)

    # -----------------------------
    # Fleet triage
    # -----------------------------
    add_section_divider(prs, footer_left, footer_right, "Fleet triage", "Decision view + recommended actions")

    slide = slide_blank(prs, "Fleet decision view", "Risk distribution + decision trade-off", footer_left, footer_right)
    hist = paths.figures / "fleet_risk_hist_latest.png"
    scat = paths.figures / "fleet_cost_vs_risk_latest.png"

    if hist.exists():
        add_media_card(slide, hist, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Fleet risk distribution (latest)")
        used_figs.add(hist.name)
    else:
        missing.append("outputs/figures/fleet_risk_hist_latest.png")

    if scat.exists():
        add_media_card(slide, scat, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Cost vs risk (latest)")
        used_figs.add(scat.name)
    else:
        missing.append("outputs/figures/fleet_cost_vs_risk_latest.png")

    add_interpretation_panel(
        slide,
        [
            "Histogram: look for a small high-risk tail (few vehicles drive most of the risk).",
            "Cost vs risk: upper-right suggests high priority (high risk and high expected cost impact).",
            "This view motivates why prioritization beats ‘first come, first served’.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Triage summary
    slide = slide_blank(prs, "Triage summary", "How many vehicles are actionable today?", footer_left, footer_right)
    triage_sum = safe_read_json(paths.tables / "triage_summary_latest.json")
    if isinstance(triage_sum, dict):
        items = [
            ("Fleet size", str(triage_sum.get("n_vehicles", "?"))),
            ("Actionable", str(triage_sum.get("n_actionable_ok", "?"))),
            ("Abstain", str(triage_sum.get("n_abstain", "?"))),
            ("Queue (K)", str(triage_sum.get("k_queue", "?"))),
        ]
        add_kpi_tiles(slide, items, MARGIN_L, BODY_TOP, full_w, Inches(1.45))
        add_interpretation_panel(
            slide,
            [
                "Actionable = passes guardrails; candidates for automatic ranking and scheduling.",
                "Abstain = OOD/low confidence; should be reviewed manually or require more data.",
                "Queue (K) maps to daily/weekly service capacity.",
            ],
            left=MARGIN_L, top=BODY_TOP + Inches(1.65), width=full_w, height=BODY_H - Inches(1.65)
        )
        used_tables.add("triage_summary_latest.json")
    else:
        missing.append("outputs/tables/triage_summary_latest.json")
        add_bullets(slide, ["Missing triage_summary_latest.json. Run: .\\scripts\\run.ps1 report"], MARGIN_L, BODY_TOP, full_w, BODY_H)

    # Mean risk by route
    slide = slide_blank(prs, "Mean risk by route", "Where risk concentrates operationally", footer_left, footer_right)
    mrr = safe_read_csv(paths.tables / "mean_risk_by_route_latest.csv")
    if mrr is not None:
        img, _ = render_table_png(mrr, deck_assets / "mean_risk_by_route.png", title="Mean risk by route (latest)", max_rows=16, max_cols=12)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, full_w, MEDIA_H, caption="Targeted planning by route")
        used_tables.add("mean_risk_by_route_latest.csv")
    else:
        missing.append("outputs/tables/mean_risk_by_route_latest.csv")
        add_bullets(slide, ["Missing mean_risk_by_route_latest.csv. Run: .\\scripts\\run.ps1 report"], MARGIN_L, BODY_TOP, full_w, MEDIA_H)

    add_interpretation_panel(
        slide,
        [
            "Higher mean risk routes may indicate harsher operating conditions or specific usage profiles.",
            "This supports operational mitigation: route-specific inspections, spares, or preventive maintenance.",
            "Treat as a planning layer (aggregate view), not a substitute for per-vehicle ranking.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Service queue
    slide = slide_blank(prs, "Recommended service queue", "Top-K prioritized vehicles", footer_left, footer_right)
    sq = safe_read_csv(paths.tables / "service_queue_latest.csv")
    if sq is not None:
        img, _ = render_table_png(sq, deck_assets / "service_queue.png", title="Service queue (latest)", max_rows=16, max_cols=12)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, full_w, MEDIA_H, caption="Ranked by cost impact (latest snapshot)")
        used_tables.add("service_queue_latest.csv")
    else:
        missing.append("outputs/tables/service_queue_latest.csv")
        add_bullets(slide, ["Missing service queue. Run: .\\scripts\\run.ps1 report"], MARGIN_L, BODY_TOP, full_w, MEDIA_H)

    add_interpretation_panel(
        slide,
        [
            "Rows are ordered by the chosen ranking objective (e.g., cost impact).",
            "Top of the list = highest expected value of servicing now (given risk + cost).",
            "Use as an input to scheduling; final decision can consider constraints (parts, location, etc.).",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Abstentions
    slide = slide_blank(prs, "Guardrails", "Abstentions (OOD / low confidence)", footer_left, footer_right)
    ab = safe_read_csv(paths.tables / "abstained_vehicles_latest.csv")
    if ab is not None:
        img, _ = render_table_png(ab, deck_assets / "abstentions.png", title="Abstained vehicles (latest)", max_rows=16, max_cols=12)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, full_w, MEDIA_H, caption="Flagged for manual review")
        used_tables.add("abstained_vehicles_latest.csv")
    else:
        missing.append("outputs/tables/abstained_vehicles_latest.csv")
        add_bullets(slide, ["Missing abstentions. Run: .\\scripts\\run.ps1 report"], MARGIN_L, BODY_TOP, full_w, MEDIA_H)

    add_interpretation_panel(
        slide,
        [
            "Abstentions are not ‘low risk’; they are ‘low trust’.",
            "Common causes: out-of-distribution features, missing signals, or probability uncertainty.",
            "Operationally: route these cases to technicians/engineers for targeted diagnosis.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Similar-case evidence
    slide = slide_blank(prs, "Similar-case evidence", "Nearest-neighbor retrieval for transparency", footer_left, footer_right)
    idx = safe_read_csv(paths.tables / "similar_cases_index_latest.csv")
    top_case = find_first([paths.tables / "similar_cases_latest_01.csv", paths.tables / "similar_cases_latest_1.csv"])

    if idx is not None:
        img, _ = render_table_png(idx, deck_assets / "similar_index.png", title="Similar-cases index", max_rows=12, max_cols=8)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, col_w, MEDIA_H, caption="Index")
        used_tables.add("similar_cases_index_latest.csv")
    else:
        missing.append("outputs/tables/similar_cases_index_latest.csv")

    if top_case is not None:
        tdf = safe_read_csv(top_case)
        if tdf is not None and not tdf.empty:
            img, _ = render_table_png(tdf, deck_assets / "similar_case_01.png", title=f"Example evidence ({top_case.name})", max_rows=12, max_cols=10)
            add_media_card(slide, img, MARGIN_L + col_w + GAP, BODY_TOP, col_w, MEDIA_H, caption="Example")
            used_tables.add(top_case.name)

    add_interpretation_panel(
        slide,
        [
            "Index lists nearest historical cases for a queried vehicle (similar signals/usage profile).",
            "Use evidence to explain the ranking: ‘this looks like past failures of type X’.",
            "This improves trust and helps engineers validate whether the recommendation makes sense.",
        ],
        left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
    )

    # Key takeaways (Fleet triage)
    slide = slide_blank(prs, "Key takeaways", "Fleet triage", footer_left, footer_right)
    bullets = [
        "Triage converts model outputs into operational actions: queue, abstentions, and explanations.",
    ]
    if isinstance(triage_sum, dict):
        bullets.append(f"Actionable vs Abstain: {triage_sum.get('n_actionable_ok','?')} vs {triage_sum.get('n_abstain','?')}.")
        bullets.append("Interpretation: abstentions represent uncertainty/OOD and should be manually reviewed.")
    else:
        bullets.append("Triage summary missing → rerun report to regenerate.")
    bullets.extend(
        [
            "Risk distribution often shows a high-risk tail → prioritization yields high ROI.",
            "Service queue is a ranked shortlist; similar-case evidence supports transparency and debugging.",
        ]
    )
    add_bullets(slide, bullets, MARGIN_L, BODY_TOP, full_w, BODY_H, font_size=18)

    # -----------------------------
    # Appendix
    # -----------------------------
    add_section_divider(prs, footer_left, footer_right, "Appendix", "Additional figures and tables (auto)")

    # Additional figures (2x2 grid)
    all_figs = sorted([p for p in paths.figures.glob("*.png") if p.is_file()], key=lambda p: p.name.lower())
    remaining_figs = [p for p in all_figs if p.name not in used_figs]
    remaining_figs = remaining_figs[: max(0, args.max_appendix_figures)]

    if remaining_figs:
        chunks = [remaining_figs[i:i + 4] for i in range(0, len(remaining_figs), 4)]
        grid_w = (full_w - GAP) / 2
        grid_h = (BODY_H - GAP) / 2
        for si, chunk in enumerate(chunks, start=1):
            slide = slide_blank(prs, "Appendix: Additional figures", f"{si}/{len(chunks)}", footer_left, footer_right)
            positions = [
                (MARGIN_L, BODY_TOP),
                (MARGIN_L + grid_w + GAP, BODY_TOP),
                (MARGIN_L, BODY_TOP + grid_h + GAP),
                (MARGIN_L + grid_w + GAP, BODY_TOP + grid_h + GAP),
            ]
            for (x, y), fp in zip(positions, chunk):
                add_media_card(slide, fp, x, y, grid_w, grid_h, caption=fp.name)

    # Additional tables (one per slide)
    all_tables = sorted([p for p in paths.tables.glob("*.csv") if p.is_file() and is_table_candidate(p)], key=lambda p: p.name.lower())
    remaining_tables = [p for p in all_tables if p.name not in used_tables]
    remaining_tables = remaining_tables[: max(0, args.max_appendix_tables)]

    for p in remaining_tables:
        df = safe_read_csv(p)
        if df is None or df.empty:
            continue
        slide = slide_blank(prs, "Appendix: Additional table", p.name, footer_left, footer_right)
        img, _ = render_table_png(df, deck_assets / f"tbl_{re.sub(r'[^a-zA-Z0-9]+','_',p.stem)}.png", title=p.name, max_rows=18, max_cols=12)
        add_media_card(slide, img, MARGIN_L, BODY_TOP, full_w, MEDIA_H, caption=p.name)
        add_interpretation_panel(
            slide,
            [
                "This table is auto-included from outputs/tables.",
                "Use it as supporting detail; keep the main story in the earlier slides.",
            ],
            left=MARGIN_L, top=note_y, width=full_w, height=NOTE_H
        )

    # Save
    out_path = Path(args.out)
    if not out_path.is_absolute():
        out_path = root / out_path
    ensure_dir(out_path.parent)
    prs.save(str(out_path))

    print(f"Deck written: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    if missing:
        missing_u = sorted(set(missing))
        print(f"Missing artifacts (non-fatal): {len(missing_u)}")
        for m in missing_u[:12]:
            print(f"  - {m}")


if __name__ == "__main__":
    main()